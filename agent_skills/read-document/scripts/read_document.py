"""
文档读取技能 - 统一读取多种文档格式
支持: PDF, Word, Excel, PowerPoint, 文本文件
"""
import argparse
import json
import logging
import sys
from pathlib import Path
from typing import Dict, Any, Optional, List

# 配置日志输出到 stderr
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    stream=sys.stderr
)
logger = logging.getLogger(__name__)


def read_pdf_content(file_path: str, page: Optional[int] = None) -> Dict[str, Any]:
    """读取PDF文件内容
    
    Args:
        file_path: PDF文件路径
        page: 指定页码（可选）
        
    Returns:
        包含PDF内容的字典
    """
    try:
        import fitz  # PyMuPDF
    except ImportError:
        return {
            "success": False,
            "error": "缺少依赖库 PyMuPDF，请运行: pip install PyMuPDF"
        }
    
    try:
        logger.info(f"开始读取PDF文件: {file_path}")
        doc = fitz.open(file_path)
        total_pages = len(doc)
        
        pages_content = []
        full_text = ""
        
        if page is not None:
            # 读取指定页
            if page < 1 or page > total_pages:
                doc.close()
                return {
                    "success": False,
                    "error": f"页码超出范围，PDF共{total_pages}页，请求第{page}页"
                }
            
            page_obj = doc[page - 1]
            text = page_obj.get_text()
            pages_content.append({
                "page_number": page,
                "content": text.strip()
            })
            full_text = text.strip()
        else:
            # 读取所有页
            for page_num in range(total_pages):
                page_obj = doc[page_num]
                text = page_obj.get_text()
                
                pages_content.append({
                    "page_number": page_num + 1,
                    "content": text.strip()
                })
                
                full_text += text + "\n\n"
        
        doc.close()
        
        logger.info(f"成功读取PDF文件，共{len(pages_content)}页")
        return {
            "success": True,
            "file_type": "pdf",
            "total_pages": total_pages,
            "pages": pages_content,
            "full_text": full_text.strip()
        }
    except fitz.FileDataError as e:
        logger.error(f"PDF文件数据错误: {file_path}, 错误: {str(e)}")
        return {
            "success": False,
            "error": f"PDF文件损坏或格式错误: {str(e)}"
        }
    except Exception as e:
        logger.error(f"读取PDF文件失败: {file_path}, 错误: {str(e)}")
        return {
            "success": False,
            "error": f"读取PDF文件失败: {str(e)}"
        }


def read_word_content(file_path: str) -> Dict[str, Any]:
    """读取Word文件内容
    
    Args:
        file_path: Word文件路径
        
    Returns:
        包含Word内容的字典
    """
    try:
        from docx import Document
    except ImportError:
        return {
            "success": False,
            "error": "缺少依赖库 python-docx，请运行: pip install python-docx"
        }
    
    try:
        logger.info(f"开始读取Word文件: {file_path}")
        doc = Document(file_path)
        paragraphs = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)
        
        # 读取表格
        tables_content = []
        for table in doc.tables:
            table_data = []
            for row in table.rows:
                row_data = [cell.text for cell in row.cells]
                table_data.append(row_data)
            tables_content.append(table_data)
        
        logger.info(f"成功读取Word文件，共{len(paragraphs)}段落，{len(tables_content)}个表格")
        return {
            "success": True,
            "file_type": "docx",
            "paragraphs": paragraphs,
            "tables": tables_content,
            "full_text": "\n".join(paragraphs)
        }
    except Exception as e:
        logger.error(f"读取Word文件失败: {file_path}, 错误: {str(e)}")
        return {
            "success": False,
            "error": f"读取Word文件失败: {str(e)}"
        }


def read_excel_content(file_path: str, sheet_name: Optional[str] = None) -> Dict[str, Any]:
    """读取Excel文件内容
    
    Args:
        file_path: Excel文件路径
        sheet_name: 指定工作表名称（可选）
        
    Returns:
        包含Excel内容的字典
    """
    try:
        import pandas as pd
    except ImportError:
        return {
            "success": False,
            "error": "缺少依赖库 pandas，请运行: pip install pandas openpyxl"
        }
    
    try:
        logger.info(f"开始读取Excel文件: {file_path}")
        
        # 读取所有工作表名称
        excel_file = pd.ExcelFile(file_path)
        
        if sheet_name:
            # 读取指定工作表
            if sheet_name not in excel_file.sheet_names:
                return {
                    "success": False,
                    "error": f"工作表 '{sheet_name}' 不存在，可用工作表: {excel_file.sheet_names}"
                }
            
            df = pd.read_excel(file_path, sheet_name=sheet_name)
            return {
                "success": True,
                "file_type": "excel",
                "sheet_name": sheet_name,
                "data": df.to_dict(orient='records'),
                "columns": df.columns.tolist(),
                "shape": df.shape
            }
        else:
            # 读取所有工作表
            all_sheets = {}
            for sheet in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet)
                all_sheets[sheet] = {
                    "data": df.to_dict(orient='records'),
                    "columns": df.columns.tolist(),
                    "shape": df.shape
                }
            
            logger.info(f"成功读取Excel文件，共{len(excel_file.sheet_names)}个工作表: {excel_file.sheet_names}")
            return {
                "success": True,
                "file_type": "excel",
                "sheet_names": excel_file.sheet_names,
                "sheets": all_sheets
            }
    except Exception as e:
        logger.error(f"读取Excel文件失败: {file_path}, 错误: {str(e)}")
        return {
            "success": False,
            "error": f"读取Excel文件失败: {str(e)}"
        }


def read_pptx_content(file_path: str) -> Dict[str, Any]:
    """读取PowerPoint文件内容
    
    Args:
        file_path: PowerPoint文件路径
        
    Returns:
        包含PPT内容的字典
    """
    try:
        from pptx import Presentation
    except ImportError:
        return {
            "success": False,
            "error": "缺少依赖库 python-pptx，请运行: pip install python-pptx"
        }
    
    try:
        logger.info(f"开始读取PPTX文件: {file_path}")
        prs = Presentation(file_path)
        slides_content = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            slides_content.append({
                "slide_number": slide_num,
                "content": "\n".join(slide_text)
            })
        
        logger.info(f"成功读取PPTX文件，共{len(prs.slides)}张幻灯片")
        return {
            "success": True,
            "file_type": "pptx",
            "total_slides": len(prs.slides),
            "slides": slides_content,
            "full_text": "\n\n".join([s["content"] for s in slides_content])
        }
    except Exception as e:
        logger.error(f"读取PPTX文件失败: {file_path}, 错误: {str(e)}")
        return {
            "success": False,
            "error": f"读取PPTX文件失败: {str(e)}"
        }


def read_text_file(file_path: str) -> Dict[str, Any]:
    """读取纯文本文件
    
    Args:
        file_path: 文本文件路径
        
    Returns:
        包含文本内容的字典
    """
    try:
        logger.info(f"开始读取文本文件: {file_path}")
        
        # 尝试多种编码
        encodings = ['utf-8', 'gbk', 'gb2312', 'latin1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    content = f.read()
                
                logger.info(f"成功读取文本文件（编码: {encoding}），共{len(content)}字符")
                return {
                    "success": True,
                    "file_type": "text",
                    "encoding": encoding,
                    "content": content,
                    "lines": content.splitlines(),
                    "line_count": len(content.splitlines())
                }
            except UnicodeDecodeError:
                continue
        
        return {
            "success": False,
            "error": "无法解码文件，尝试了编码: " + ", ".join(encodings)
        }
    except Exception as e:
        logger.error(f"读取文本文件失败: {file_path}, 错误: {str(e)}")
        return {
            "success": False,
            "error": f"读取文本文件失败: {str(e)}"
        }


def read_document(file_path: str, sheet_name: Optional[str] = None, page: Optional[int] = None) -> Dict[str, Any]:
    """统一文档读取接口
    
    Args:
        file_path: 文档路径
        sheet_name: Excel工作表名称（可选）
        page: PDF页码（可选）
        
    Returns:
        包含文档内容的字典
    """
    path = Path(file_path)
    
    if not path.exists():
        return {
            "success": False,
            "error": f"文件不存在: {file_path}"
        }
    
    ext = path.suffix.lower()
    
    # 根据文件扩展名选择读取方法
    if ext == '.pdf':
        return read_pdf_content(file_path, page)
    elif ext == '.docx':
        return read_word_content(file_path)
    elif ext in ['.xlsx', '.xls']:
        return read_excel_content(file_path, sheet_name)
    elif ext == '.pptx':
        return read_pptx_content(file_path)
    else:
        # 默认作为文本文件处理
        return read_text_file(file_path)


def get_supported_formats() -> Dict[str, Any]:
    """获取支持的文件格式列表"""
    return {
        "success": True,
        "supported_formats": {
            "pdf": {
                "extensions": [".pdf"],
                "description": "PDF文档",
                "library": "PyMuPDF"
            },
            "word": {
                "extensions": [".docx"],
                "description": "Word文档",
                "library": "python-docx"
            },
            "excel": {
                "extensions": [".xlsx", ".xls"],
                "description": "Excel表格",
                "library": "pandas + openpyxl"
            },
            "powerpoint": {
                "extensions": [".pptx"],
                "description": "PowerPoint演示文稿",
                "library": "python-pptx"
            },
            "text": {
                "extensions": [".txt", ".md", ".csv", ".json", ".xml", ".html", ".css", ".js", ".py"],
                "description": "文本文件",
                "library": "内置"
            }
        },
        "note": "根据文件扩展名自动识别并选择相应的解析方法"
    }


def main():
    parser = argparse.ArgumentParser(description="文档读取技能 - 统一读取多种文档格式")
    parser.add_argument("--file", required=True, help="文档文件路径")
    parser.add_argument("--sheet", default=None, help="Excel工作表名称（可选）")
    parser.add_argument("--page", type=int, default=None, help="PDF页码（可选）")
    parser.add_argument("--formats", action="store_true", help="列出支持的文件格式")
    
    args = parser.parse_args()
    
    if args.formats:
        result = get_supported_formats()
    else:
        result = read_document(args.file, args.sheet, args.page)
    
    # 输出 JSON 结果到 stdout
    print(json.dumps(result, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
